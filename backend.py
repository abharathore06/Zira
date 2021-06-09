import pyttsx3                              #library that supports converting text-to-speech
import datetime                             #for getting the live time
import speech_recognition as sr             #library that will recognize user takeCommand
import wikipedia                            #it is used to access wikipedia
import webbrowser                           #it is used to access browser things
import os                                   #it is used to access our local computer files
import smtplib                              #to get access to mail server      
import random                               #to get random songs
import time                                 #to get current time
import pyjokes                              #to generate random jokes
import subprocess                           #to process tasks like shut down, logout and more
import ctypes                               #allows us to change background image
import pyautogui                            #allows us to take a screenshot
import requests                             #allows us to send HTTP requests
import wolframalpha as wa                   #allows us to access question and solve them
import cv2                                  #allows us to open webcam
import pywhatkit as kit                     #allows to search and  play
import winshell                             #to access windows functionality
import gmapFinder                           #used to access maps for directions
import xlsxwriter                           #used to write in excel
from openpyxl import *                      #used to load excel data
from pptx import Presentation               #used to write, save in powerpoint
import googletrans                          #used to translate in various languages
from googletrans import Translator, LANGUAGES
import glob                                 #used to retrieve files/pathnames matching a specified pattern.       
import psutil                               #used here for battery information
from datetime import date                   #used to fetch the current date
import wmi                                  #used to get system information
import screen_brightness_control as sbc     #used to control screen brightness
import keyboard

engine = pyttsx3.init('sapi5')                  #speech api helps in recognition of voice        
voices=engine.getProperty('voices')             #get details of voice
engine.setProperty('voice', voices[1].id)       #selecting the voice
sys_sound = 0                                   


def speak(audio):
    engine.say(audio)               #function say, converts text to speech
    engine.runAndWait()             #makes the voice audible


def wishme():
    hour=int(datetime.datetime.now().hour)              #gives the live time
    if(hour>=0 and hour<12):
        speak("Good Morning")
    elif(hour>=12 and hour<18):
        speak("Good afternoon!")
    elif(hour>=18 and hour<21):
        speak("Good evening!")
    else:
        speak("Good night!")
    speak("I am Zira, I am your personal assistant. Please tell me how may I help you?")
    

def takeCommand():
    #It takes microphone takeCommand from the user and returns string output

    r = sr.Recognizer()                         #recognizer is a in-built class
    with sr.Microphone() as source:             #takeCommand will store in source
        print("Listening...")
        r.pause_threshold = 2
        audio = r.listen(source)

    try:
        print("Recognizing...")    
        query = r.recognize_google(audio, language='en-in')              #Using google for voice recognition.
        print(f"User said: {query}\n")                                   #User query will be printed.

    except Exception as e:
        # print(e)    
        print("Pardon me, Say that again please...")           #this will be printed if system doesn't get a proper voice takeCommand
        return "None"                               #None string will be returned
    return query

def sendEmail(to, content):
    server = smtplib.SMTP("smtp.gmail.com", 587)
    server.ehlo()
    server.starttls()
    server.login("aabharathore86862@gmail.com", "yashita.26.rathore")
    server.sendmail("aabharathore86862@gmail.com", to, content)
    server.close()

def get_size(bytes, suffix="B"):
    """
    Scale bytes to its proper format
    e.g:
        1253656 => '1.20MB'
        1253656678 => '1.17GB'
    """
    factor = 1024
    for unit in ["", "K", "M", "G", "T", "P"]:
        if bytes < factor:
            return f"{bytes:.2f}{unit}{suffix}"
        bytes /= factor


if __name__=="__main__" :

    wishme()

    while True:
    #if 1:
        query = takeCommand().lower()       #making it in lower case so it can match our code

        if 'tell me about yourself' in query:
            speak("I am Zira your personal assistant. I am programmed to perform minor tasks like opening websites for you ,predict time,take a photo,search wikipedia,predict weather In different cities, get top headlines of news, launch applications and more.")

        
        
        
        elif "who made you" in query or "who created you" in query:
            speak("I was built by Abha")
            print("I was built by Abha")

        


        # speak("What's your name, Human?")
        # name ='Human'
        # name = takeCommand()
        # speak("Hello, " + name + '.')
        


        
        elif "how are you" in query:
            speak("I am fine, Thank you!")
            speak("How are you?")

        
        
        
        elif 'joke' in query:
            speak(pyjokes.get_joke(language='en', category = 'all'))

        
        
        
        elif "who i am" in query:
            speak("If you talk then definately your human.")
        
        
        
        
        elif "voice" in query:
            if sys_sound == 0:
                sys_sound = 1
            else:
                sys_sound = 0    
            engine.setProperty('voice', voices[sys_sound].id) 
            speak("As per your request I have changed  my voice")
 
        
        
        
        elif "directions" in query or "direction" in query:
            speak("At what point you want the directions from?")
            curr = takeCommand()
            speak("Okay, tell me where do you want to go?")
            dest = takeCommand()
            directions = gmapFinder.getDirection(curr, dest)
            speak("Here are the directions from" + curr + "to" + dest)
            webbrowser.open(directions["url"])
            #print(directions["url"])

        


        elif "date" in query:
            today = date.today()
            d = today.strftime("%B %d, %Y")
            print("date =", d)
            speak("Today's date is ")
            speak(d)
        
        


        elif 'change background' in query:
            my_wall = 'E:\\backgrounds'
            bg = os.listdir(my_wall)
            folder = random.choice(bg)
            ctypes.windll.user32.SystemParametersInfoW(20,0,folder,0)
            speak("Background changed succesfully")

        
        
        
        elif "fine" in query or "good" in query:
            speak("It's good to know that you are fine.")

        
        
        
        elif "cost of" in query:
            search_term = query
            url = "https://google.com/search?q=" + search_term
            webbrowser.get().open(url)
            speak("Here is what I found for " + search_term + " on google")   

        
        
        
        elif "take a picture" in query or "take a photo" in query:
            cap = cv2.VideoCapture(0)
            ret, frame = cap.read()
            print (ret)
            cv2.imwrite("NewPicture.jpg",frame)
            cap.release()

        
        
        
        elif "show picture" in query or 'show photo' in query:
            cv2.imshow('frame', frame)
            cv2.waitKey()
            cv2.destroyAllWindows()




        elif "open a new tab" in query or "open tab" in query:
            keyboard.press_and_release("ctrl + t")
            speak("Opened a new tab")




        elif "open a new window" in query or "open a new browser window" in query or "open window" in query:
            keyboard.press_and_release("ctrl + n")
            speak("Opened a new browser window")




        elif "reload" in query or "refresh" in query:
            keyboard.press_and_release("f5")
            speak("Reloaded the web")




        elif "full screen" in query:
            keyboard.press_and_release("f11")
            speak("Switched to the full screen window")




        elif "scroll" in query: 
            print("You want to do scroll up or scroll down?")
            speak("You want to do scroll up or scroll down?")
            comm = takeCommand()
            if "up" in comm:
                keyboard.press_and_release("shift + space")
            elif "down" in comm:
               keyboard.press_and_release("space") 



        
        elif "browsing history" in query or "history" in query or "browser history" in query or "browsed history" in query:
            keyboard.press_and_release("ctrl + h")
            speak("Opened your browsing history")




        elif "downloading history" in query or "download" in query or "download history" in query or "downloded history" in query:
            keyboard.press_and_release("ctrl + j")
            speak("Opened your download history")




        elif "open bookmark" in query : 
            keyboard.press_and_release("ctrl + shift + o")
            speak("Opened your Bookmark Manager") 




        elif "bookmark" in query : 
            keyboard.press_and_release("ctrl + d")
            speak("Bookmarked the current website")




        elif "close this tab" in query or "close the tab" in query or "close tab" in query:
            keyboard.press_and_release("ctrl + w")
            speak("Closed the tab")



        
        elif "switch the tab" in query or "switch tab" in query or "switch this tab" in query:
            print("Which tab you want to switch with?")
            speak("Which tab you want to switch with?")
            instr = takeCommand()
            if "last" in instr or "recently opened" in instr or "end" in instr or "ending" in instr:
                keyboard.press_and_release("ctrl + 9")
                speak("Switched to the last tab")
            elif "first" in instr or "firstly opened" in instr or "start" in instr or "starting" in instr:
                keyboard.press_and_release("ctrl + 1")
                speak("Switched to the first tab")
            elif "next" in instr:
                keyboard.press_and_release("ctrl + tab")
                speak("Switched to the next tab")
            elif "previous" in instr:
                keyboard.press_and_release("ctrl + shift + tab")
                speak("Switched to the previous tab")




        elif "reopen the last closed tab" in query or "reopen the tab" in query or "reopen tab" in query or "reopen last tab" in query or "reopen closed tab" in query:
            keyboard.press_and_release("ctrl + shift + t")
            speak("Reopened the last closed tab")




        elif "mute" in query or "mute volume" in query: 
            pyautogui.press('volumemute')
            speak("Volume has been muted")




        elif "raise" in query or "increase" in query or "up" in query or "high" in query or "higher" in query: 
            pyautogui.press('volumeup')
            speak("Volume has been raised")




        elif "down" in query or "decrease" in query or "lower" in query or "low" in query:
            pyautogui.press('volumedown')
            speak("Volume has been decreased")




        elif "play" in query or "pause" in query:
            pyautogui.press('playpause')




        elif "timer" in query or "countdown" in query:
            print("Enter the time in seconds: ")
            speak("Enter the time in seconds: ")
            t = takeCommand()
            while t:
                mins, secs = divmod(t, 60)
                timer = '{:02d}:{:02d}'.format(mins, secs)
                print(timer, end="\r")
                time.sleep(1)
                t -= 1
            print('Fire in the hole!! Your timer has been ended')
            speak('Fire in the hole!! Your timer has been ended')




        elif "battery" in query:
            # function returning time in hh:mm:ss
            def convertTime(seconds):
                minutes, seconds = divmod(seconds, 60)
                hours, minutes = divmod(minutes, 60)
                return "%d:%02d:%02d" % (hours, minutes, seconds)
  
            # returns a tuple
            battery = psutil.sensors_battery()
  
            print("\nBattery charged : ", battery.percent,"%")
            speak(f"The battery is {battery.percent} % charged.")
            if battery.percent<25 and battery.power_plugged == False:
                print("Your system will soon go into Power Saving mode. I would suggest that you get it charged.")
                speak("Your system will soon go into power saving mode. I would suggest that you get it charged.")
            if battery.power_plugged == True:
                print("The battery is plugged-in.")
                speak("It is plugged in.")
            elif battery.percent>30 and battery.power_plugged == False:
                print("The battery is not plugged-in.")
                speak("The battery is not plugged in.")
  
            # converting seconds to hh:mm:ss
            time_str = convertTime(battery.secsleft).replace('-',' ')
            print("Battery left : ", convertTime(battery.secsleft)) 
            speak(f"And it will last for {time_str}")


        
        
        
        elif 'game' in query:
            speak("Yeah! I have some games for you")
            print("1. Roshambo \n 2. Online games")
            speak("Which one would you like to play?")
            choice = takeCommand()
            if '1' in choice or 'roshambo' in choice:
                speak("choose among rock paper or scissor")
                user = takeCommand()
                moves=["rock", "paper", "scissor"]
                g_move=random.choice(moves)
                u_move=user
                speak("The computer chose " + g_move)
                speak("You chose " + u_move)
                #engine_speak("hi")
                if u_move==g_move:
                    speak("the match is draw")
                elif u_move== "rock" and g_move== "scissor":
                    speak("Player wins")
                elif u_move== "rock" and g_move== "paper":
                    speak("Computer wins")
                elif u_move== "paper" and g_move== "rock":
                    speak("Player wins")
                elif u_move== "paper" and g_move== "scissor":
                    speak("Computer wins")
                elif u_move== "scissor" and g_move== "paper":
                    speak("Player wins")
                elif u_move== "scissor" and g_move== "rock":
                    speak("Computer wins")
            
            elif '2' in choice or 'online games' in choice:
                webbrowser.open("https://iogames.space/")
                speak("Here are some online games for you. Checkout this")

        
        
        
        elif "song for me" in query:
            speak("Sure! Which song would you like me to play?")
            song = takeCommand()
            kit.playonyt(song)




        elif "translate" in query:
            translator = Translator()
            #print(googletrans.LANGUAGES)
            key_list = list(LANGUAGES.keys())
            val_list = list(LANGUAGES.values())

            print('What text would you like me to translate?') 
            speak('What text would you like me to translate?') 
            text = str(takeCommand())

            print("To which language I should translate?")
            speak("To which language I should translate?")
            lang = str(takeCommand())
            
            
            position = val_list.index(lang.lower())
            language = key_list[position]
            #print(language)

            result = translator.translate(text, src = "auto", dest = language)   
            print(result.text)
            speak(result.text)




        elif "system information" in query:
            c = wmi.WMI()
            my_system = c.Win32_ComputerSystem()[0]

            print(f"Manufacturer: {my_system.manufacturer}")
            speak(f"Manufacturer: {my_system.manufacturer}")
            print(f"Device Name: {my_system.name}")
            speak(f"Device Name: {my_system.name}")
            print(f"Model: {my_system.model}")
            speak(f"Model: {my_system.model}")
            print(f"System Type: {my_system.systemType}")
            speak(f"System Type: {my_system.systemType}")
            svmem = psutil.virtual_memory()
            print(f"Total Memory :{get_size(svmem.total)}")
            speak(f"Total Memory :{get_size(svmem.total)}")
            print(f"Available: {get_size(svmem.available)}")
            speak(f"Available: {get_size(svmem.available)}")
            print(f"Used: {get_size(svmem.used)}")
            speak(f"Used: {get_size(svmem.used)}")
            print(f"Percentage: {svmem.percent}%")
            speak(f"Percentage: {svmem.percent}%")




        elif "convert" in query:
            speak("Tell me what I have to convert?")
            text_c = takeCommand() 
            kit.text_to_handwriting(text_c)

        
        
        
        elif "search for" in query:
            kit.search(query)

        

        
        elif 'send a whatsapp message' in query or 'send whatsapp' in query:
            phone_dict={"abha":"+917746081400", "shreya":"+919352056977", }#use only lower case letters
            print('\nWho should I send this message to?')
            speak('Who should I send this message to?')
            name=takeCommand()
            if name in phone_dict:
                ph_no=str(phone_dict[name])
                print('\nWhat would you like it to say?')
                speak('What would you like it to say?')
                msg=takeCommand()
                content=(f'J.A.R.V.I.S. :  {msg}')
                print('\nWould you like to send this message now or at a specific time?')
                speak('Would you like to send this message now or at a specific time?')
                choice = takeCommand()
                if 'now' in choice:
                    kit.sendwhatmsg_instantly(ph_no,content)   
                elif 'time' in choice:
                    print("\nSpecify hour as in 24 Hr format")
                    speak("Specify hour as in 24 hour format")
                    h=int(takeCommand())
                    print("\nSpecify minutes")
                    speak("Specify minutes")
                    m=int(takeCommand())
                    kit.sendwhatmsg(ph_no, content, h , m)
                else:
                    print("\nMessage cancelled")
                    speak("Message cancelled")   
                            
            else:
                print(f"\nI could not find '{name}' in your contacts.")  
                speak(f"I could not find {name} in your contacts.")




        elif "make a note" in query or "write a note" in query:
            speak("What should I write dear?")
            note=takeCommand()
            file = open('Zira.txt', 'w')
            speak("Should I include date and time?")
            ans = takeCommand()
            if 'yes' in ans or 'sure' in ans:
                strTime = datetime.datetime.now().strftime("%H:%M:%S")
                file.write(strTime)
                file.write("\n")
                file.write(note)
            else:
                file.write(note)




        elif "brightness" in query or "light" in query:
            current_b = sbc.get_brightness()
            print("Your current brightness level is", current_b, "%.")
            speak(f"Your current brightness level is {current_b} % .")
            print("Would you like to adjust?")
            speak("Would you like to adjust?")
            com = takeCommand()

            if 'yes' in com or "sure" in com or "yeah" in com or "ok" in com:
                print("At what percentage you want to adjust?")
                speak("At what percentage you want to adjust?")
                per = int(takeCommand())

                if per<current_b:
                    new1 = sbc.set_brightness(per)
                    print("Your brightness is faded to ", new1, "%")
                    speak(f"Your brightness is faded to {new1} percent")
                elif per>current_b:
                    new2 = sbc.set_brightness(per)
                    print("Your brightness is increased to ", new2, "%")
                    speak(f"Your brightness is increased to {new2} percent")
                else:
                    print("That's exactly the level you're at.")
                    speak("That's exactly the level you're at.")

            else:
                print("OK, You are on the same brightness level")
                speak("OK, You are on the same brightness level")




        elif "minimize" in query:
            fw = pyautogui.getActiveWindow()
            #print(fw)
            if (fw.isMinimized==True):
                print("The window is already in a minimized mode")
                speak("The window is already in a minimized mode")
            else:
                fw.minimize()




        elif "maximize" in query:
            fw = pyautogui.getActiveWindow()
            #print(fw)
            if (fw.isMaximized==True):
                print("The window is already in a maximized mode")
                speak("The window is already in a maximized mode")
            else:
                fw.maximize()




        elif "close" in query or "destroy" in query or "stop" in query:
            print("Do you really want to quit the window?")
            speak("Do you really want to quit the window?")
            confirm = takeCommand()
            if "yes" in confirm or "sure" in confirm or "yeah" in confirm or "sure" in confirm:
                keyboard.press_and_release("alt + f4")
            else:
                pass




        elif "show notes" in query:
            speak("Showing Notes")
            file = open('Zira.txt', 'r')
            print(file.read())
            result = "You have wrote" + file.read(6)
            speak(result)




        elif "weather" in query:
            api_key="daa1b9ca3d62bc73f0e8f7b52b67b6c3"
            base_url="https://api.openweathermap.org/data/2.5/weather?"
            speak("what is the city name")
            city_name=takeCommand()
            complete_url=base_url+"appid="+api_key+"&q="+city_name
            
            response = requests.get(complete_url)
            x=response.json()       #converting json format to python
            if x["cod"]!="404":
                y=x["main"]             #y is the main key
                current_temperature = y["temp"]
                current_humidiy = y["humidity"]
                z = x["weather"]
                weather_description = z[0]["description"]
                speak(" Temperature in kelvin unit is " +
                      str(current_temperature) +
                      "\n humidity in percentage is " +
                      str(current_humidiy) +
                      "\n and its  " +
                      str(weather_description) + "outside")
                print(" Temperature in kelvin unit = " +
                      str(current_temperature) +
                      "\n humidity (in percentage) = " +
                      str(current_humidiy) +
                      "\n description = " +
                      str(weather_description))
            else:
                speak("City not found!")




        elif "html" in query:
            speak("Tell me the Name for Project: ")
            project_name=takeCommand()
            while (os.path.exists(project_name)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                project_name = takeCommand()
                continue
            
            os.mkdir(project_name)
            speak('Successfully Created. Have a look at the project!')

            os.chdir(project_name)          #changes the current working directory
            htmlContent = '<html>\n\t<head>\n\t\t<title> ' + project_name + ' </title>\n\t\t<link rel="stylesheet" type="text/css" href="style.css">\n\t</head>\n<body>\n\t<p id="label"></p>\n\t<button id="btn" onclick="showText()"> Click Me </button>\n\t<script src="script.js"></script>\n</body>\n</html>'

            htmlFile = open('index.html', 'w')
            htmlFile.write(htmlContent)
            htmlFile.close()

            cssContent = '* {\n\tmargin:0;\n\tpadding:0;\n}\nbody {\n\theight:100vh;\n\tdisplay:flex;\n\tjustify-content:center;\n\talign-items:center;\n}\n#btn {\n\twidth:200px;\n\tpadding: 20px 10px;\n\tborder-radius:5px;\n\tbackground-color:red;\n\tcolor:#fff;\n\toutline:none;border:none;\n}\np {\n\tfont-size:30px;\n}'

            cssFile = open('style.css', 'w')
            cssFile.write(cssContent)
            cssFile.close()

            jsContent = 'function showText() {\n\tdocument.getElementById("label").innerHTML="Successfully Created Project";\n\tdocument.getElementById("btn").style="background-color:green;"\n}'

            jsFile = open('script.js', 'w')
            jsFile.write(jsContent)
            jsFile.close()

            os.startfile("index.html")

            
        

        elif "powerpoint" in query:
            speak("Tell me the name for the project")
            name = takeCommand()
            filename = name + ".pptx"
            while (os.path.exists(filename)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                name = takeCommand()
                filename = name + ".pptx"
                continue
            prs = Presentation()
            lyt = prs.slide_layouts[0]                  #choose a slide layout
            slide = prs.slides.add_slide(lyt)           #adding a slide
            title = slide.shapes.title                  #adding title
            subtitle = slide.placeholders[1]            #placeholder for subtitle
            speak("What should i write in the title section?")
            title.text = takeCommand()
            if 'nothing' in title.text or 'empty' in title.text or "leave" in title.text or "blank" in title.text or "don't want" in title.text:
                title.text = " "
            speak("Okay, now tell me for the subtitle section?")
            subtitle.text = takeCommand()
            if 'nothing' in subtitle.text or 'empty' in subtitle.text or "leave" in subtitle.text or "blank" in subtitle.text or "don't want" in subtitle.text:
                subtitle.text = " "
            prs.save(filename)
            speak("Successfully created!")
            os.startfile(filename)




        elif "excel" in query:
            row=0
            col=0
            speak("Tell me the name for the sheet?")
            name = takeCommand()
            filename = name + ".xlsx"
            while (os.path.exists(filename)==True):
                speak("Project already exists")
                speak("Tell me the another name")
                name = takeCommand()
                filename = name + ".xlsx"
                continue
            workbook= xlsxwriter.Workbook(filename)
            worksheet = workbook.add_worksheet()
            speak("What would i write in the sheet?")
            content = takeCommand()
            if "blank" in content or "empty" in content or "leave" in content or "nothing" in content or "don't want" in content:
                content = " "
            f_content = content.split()
            length = len(f_content)
            for i in range(length):
                worksheet.write(row, col,f_content[i])
                col += 1
            speak("Successfuly created. Have a look at your file")
            workbook.close()
            os.startfile(filename)




        elif 'wikipedia' in query:
            speak("Searching wikipedia.....")
            query = query.replace("wikipedia", " ")
            results = wikipedia.summary(query, sentences = 1)
            speak("According to Wikipedia")
            print(results)
            speak(results)
        



        elif 'open youtube' in query:
            speak("Opening Youtube.....")
            webbrowser.open("youtube.com")
            



        elif 'open gmail' in query:
            speak("Opening Gmail.....")
            webbrowser.open("gmail.com")



        
        elif 'open google' in query:
            speak('Opening Google')
            url='google.com'
            webbrowser.register('chrome',None, webbrowser.BackgroundBrowser("C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"))
            webbrowser.get('chrome').open(url)




        elif 'open windy' in query:
            speak("Opening Windy.....")
            webbrowser.open("windy.com")




        elif 'close the window' in query:
            os.close()

        


        elif 'play music' in query:
            music_dir = 'D:\\Zira\\fav'
            print("\nWould you like to pick a choice or would you like me to surprise you?")
            speak("Would you like to pick a choice or would you like me to surprise you?")
            choice=takeCommand()           
            if 'choice' in choice or 'first' in choice:
                #do this
                print("\nName your choice please..")
                speak("name your choice please")
                pick=takeCommand() 
                songs = glob.glob(music_dir+'/*.mp3')
                index=0
                for i in songs:
                    if pick in i.lower():
                       
                        print('\nPlaying Music..')
                        speak('Playing music')                          #filters out all the mp3 files 
                        os.startfile(os.path.join(music_dir,songs[index]))
                        break
                    elif index==(len(songs)-1):
                        print("\nSong not found")
                        speak("song not found")
                    else:
                        index=index+1
                        continue
            
            elif 'surprise' in choice or 'second' in choice:
                #do that    
                songs = glob.glob(music_dir+'/*.mp3')
                print('\nPlaying Music..')
                speak('Playing music')                      
                os.startfile(os.path.join(music_dir,random.choice(songs)))




        elif 'time' in query:
            strTime = datetime.datetime.now().strftime("%H:%M:%S")
            speak(strTime)




        elif 'capture' in query or 'screenshot' in query:
            try:
                myScreenshot = pyautogui.screenshot()
                myScreenshot.save('C:\\Users\\Asus\\OneDrive\\Pictures\\Screenshots\\ss.png') 
                speak("The screenshot has been taken successfully")
            except Exception as e:
                print(e)
                speak("Sorry, there are some issues taking a screenshot")




        elif 'open code' in query:
            speak("Opening Code.....")
            vs_path = "C:\\Users\\Asus\\AppData\\Local\\Programs\\Microsoft VS Code\\Code.exe"
            os.startfile(vs_path)




        elif 'open gallery' in query or "open photos" in query:
            speak("Opening Gallery.....") 
            mv_path = "E:\\Yashuu"
            os.startfile(mv_path)




        elif 'send an email' in query or 'send mail' in query:
            speak('yes I am working on it')
            contacts={"shreya":"shreya.awesome2040@gmail.com", "atharv" : "aaryan.a.diwankar@gmail.com"}
            print("\nWho am I writing to?")
            speak("Who am I writing to?")
            name=takeCommand()
            if name in contacts:
                mail_id=str(contacts[name])
                try:
                    print('\nWriting Mail..')
                    speak("Writing Mail")
                    print("\nWhat is the subject?")
                    speak("What is the subject?")
                    sub = takeCommand()
                    print("\nWhat should I say?")
                    speak("What should I say?")
                    msg = takeCommand()
                    content = 'Subject: {}\n\n{}'.format(sub, msg)
                    
                    print('\nTransmitting Message')
                    speak('Transmitting message')
                    sendEmail(mail_id, content)
                    print("\nTransmission Successful")
                    speak("Transmission successful")
                except Exception as e:
                    print('\nTransmission Failed\n')
                    speak('Transmission failed')
                    print(e)
                    speak(e)
            else:
                print(f"\nI could not find '{name}' in your contacts")  
                speak(f"I could not find {name} in your contacts") 






        elif 'news' in query:
            webbrowser.open("https://timesofindia.indiatimes.com/home/headlines")
            speak('Here are some headlines from the Times of India,Happy reading')
            time.sleep(5)                   #to wait till 5 seconds before the next execution




        elif "calculate" in query:
            app_id = "4WY52G-U6687A88X2"
            client = wa.Client(app_id)
            indx = query.lower().split().index('calculate')
            query = query.split()[indx + 1:]
            res = client.query(' '.join(query))
            answer = next(res.results).text
            print("The answer is " + answer)
            speak("The answer is " + answer)




        elif 'empty recycle bin' in query:
            winshell.recycle_bin().empty(confirm = False, show_progress = False, sound = True)
            speak("Recycle Bin Recycled")




        elif "where is" in query or "what is" in query or "who is" in query: 
            app_id = "4WY52G-U6687A88X2"
            client = wa.Client(app_id)
            res = client.query(query)
            answer = next(res.results).text
            try:
                print (next(res.results).text)
                speak (next(res.results).text)
            except StopIteration:
                print ("No results")




        elif "log off" in query or "shut down" in query:
            speak("Ok , your pc will lshut down in 10 sec make sure you exit from all applications")
            subprocess.call(["shutdown", "/l"])




        elif "restart" in query:
            speak("Putting the pc to restart")
            subprocess.call(["shutdown", "/r"])




        elif "hibernate" in query or "sleep" in query:
            speak("Hibernating")
            subprocess.call("shutdown / h")




        elif "quit" in query:
            speak("ok good Bye. Catch me later for more exicting things")
            exit()



 